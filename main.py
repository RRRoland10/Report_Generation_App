import argparse
from pptx import Presentation
import json

def load_configuration(config_file):
    with open(config_file, 'r') as file:
        config_data = json.load(file)
    return config_data

def generate_pptx_report(config):
    presentation = Presentation()

    # Add a title slide
    title_slide_layout = presentation.slide_layouts[0]
    title_slide = presentation.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = config.get('title', 'Report Title')
    subtitle.text = config.get('subtitle', 'Subtitle')

    # Add content slides
    content_slide_layout = presentation.slide_layouts[1]

    for content in config.get('content', []):
        content_slide = presentation.slides.add_slide(content_slide_layout)
        title = content_slide.shapes.title
        content_textbox = content_slide.placeholders[1]

        title.text = content.get('title', 'Content Title')
        content_textbox.text = content.get('content', 'Content')

    # Save the presentation
    presentation.save(config.get('output_file', 'output.pptx'))

def main():
    parser = argparse.ArgumentParser(description='Generate PowerPoint report from a configuration file.')
    parser.add_argument('config_file', help='Path to the configuration file')
    args = parser.parse_args()

    config = load_configuration(args.config_file)
    generate_pptx_report(config)

if __name__ == "__main__":
    main()
